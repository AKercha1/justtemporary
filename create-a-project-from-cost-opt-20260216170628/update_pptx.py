#!/usr/bin/env python3
"""
Update kickoff-deck.pptx with meeting information from markdown files.
Adds new slides with meeting summary and action items using existing slide styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the existing presentation
prs = Presentation('artifacts/kickoff-deck.pptx')

# Get the slide dimensions and styling from the first existing slide
if len(prs.slides) > 0:
    reference_slide = prs.slides[0]
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
else:
    blank_slide_layout = prs.slide_layouts[6]

# Create Slide 1: Meeting Summary
slide1 = prs.slides.add_slide(blank_slide_layout)

# Add title
title_box1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame1 = title_box1.text_frame
title_frame1.text = "Meeting Summary"
title_para1 = title_frame1.paragraphs[0]
title_para1.font.size = Pt(54)
title_para1.font.bold = True
title_para1.font.color.rgb = RGBColor(0, 51, 102)

# Add content
content_box1 = slide1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
content_frame1 = content_box1.text_frame
content_frame1.word_wrap = True

# Add meeting details
p = content_frame1.paragraphs[0]
p.text = "Audio Connection Verification"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(64, 64, 64)
p.space_after = Pt(18)

p = content_frame1.add_paragraph()
p.text = "Conducted by: Andrey Kerchin"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(64, 64, 64)
p.space_after = Pt(12)
p.level = 0

p = content_frame1.add_paragraph()
p.text = "Purpose: Verify audio functionality and system connectivity"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(64, 64, 64)
p.space_after = Pt(12)
p.level = 0

p = content_frame1.add_paragraph()
p.text = "Status: Connection checks and numerical counts performed successfully"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(64, 64, 64)
p.space_after = Pt(12)
p.level = 0

# Create Slide 2: Recommendations & Action Items
slide2 = prs.slides.add_slide(blank_slide_layout)

# Add title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame2 = title_box2.text_frame
title_frame2.text = "Recommendations & Action Items"
title_para2 = title_frame2.paragraphs[0]
title_para2.font.size = Pt(54)
title_para2.font.bold = True
title_para2.font.color.rgb = RGBColor(0, 51, 102)

# Add content
content_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
content_frame2 = content_box2.text_frame
content_frame2.word_wrap = True

# Add recommendations section
p = content_frame2.paragraphs[0]
p.text = "Recommendations"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)
p.space_after = Pt(12)

p = content_frame2.add_paragraph()
p.text = "No specific recommendations were formally agreed upon in this session"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(64, 64, 64)
p.space_after = Pt(18)
p.level = 0

# Add action items section
p = content_frame2.add_paragraph()
p.text = "Action Items"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)
p.space_after = Pt(12)

p = content_frame2.add_paragraph()
p.text = "No formal action items were agreed upon"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(64, 64, 64)
p.space_after = Pt(12)
p.level = 0

# Save the updated presentation
prs.save('artifacts/kickoff-deck.pptx')
print("✓ Successfully updated artifacts/kickoff-deck.pptx")
print(f"✓ Total slides: {len(prs.slides)}")
print("✓ Added 2 new slides with meeting summary and recommendations")
