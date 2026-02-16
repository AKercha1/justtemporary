#!/usr/bin/env python3
"""
Script to add meeting summary and action items slides to kickoff-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Load the presentation
prs = Presentation('artifacts/kickoff-deck.pptx')

print(f"Original total slides: {len(prs.slides)}")

# Remove the last 2 slides if they exist (in case of re-run)
while len(prs.slides) > 13:
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

print(f"Reset to original: {len(prs.slides)} slides")

# Find the appropriate layout for content slides
slide_layout = None
for layout in prs.slide_layouts:
    if 'Content Slide_with Eyebrow' in layout.name:
        slide_layout = layout
        break

if not slide_layout:
    slide_layout = prs.slide_layouts[1]

print(f"Using layout: {slide_layout.name}")

# Parse the meeting summary
with open('artifacts/meeting-summary.md', 'r') as f:
    summary_text = f.read()

# Extract specific fields from the markdown
conducted_by = "Andrey Kerchin"
purpose = "Technical test - connection and meeting test verification"
key_discussions = [
    "Audio checks performed (111, 123123 проверка связи, Раз 23 раз, 23)",
    "Meeting test confirmation (Это проверка митинга)",
    "Name mentioned: Валентика t",
    "Brief unscripted test call to verify audio and meeting functionality"
]

recommendations = []
action_items = []

# --- SLIDE 1: Meeting Summary ---
slide1 = prs.slides.add_slide(slide_layout)

# Find and set title placeholder
title_shape = None
content_shape = None

for shape in slide1.placeholders:
    if shape.placeholder_format.type == 1:  # Title
        title_shape = shape
    elif shape.placeholder_format.type == 2:  # Body
        content_shape = shape

if title_shape:
    title_shape.text = "Meeting Summary"

# Set body content
if content_shape and hasattr(content_shape, 'text_frame'):
    tf = content_shape.text_frame
    tf.clear()

    # Create paragraphs for each piece of information
    p = tf.paragraphs[0]
    p.text = f"Conducted By: {conducted_by}"
    p.level = 0
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = f"Purpose: {purpose}"
    p.level = 0
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Key Discussion Points:"
    p.level = 0
    p.font.size = Pt(14)
    p.font.bold = True

    for item in key_discussions:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(13)

# --- SLIDE 2: Recommendations & Action Items ---
slide2 = prs.slides.add_slide(slide_layout)

# Find and set title placeholder
title_shape2 = None
content_shape2 = None

for shape in slide2.placeholders:
    if shape.placeholder_format.type == 1:  # Title
        title_shape2 = shape
    elif shape.placeholder_format.type == 2:  # Body
        content_shape2 = shape

if title_shape2:
    title_shape2.text = "Recommendations & Action Items"

# Set body content
if content_shape2 and hasattr(content_shape2, 'text_frame'):
    tf2 = content_shape2.text_frame
    tf2.clear()

    p = tf2.paragraphs[0]
    p.text = "Agreed Recommendations:"
    p.level = 0
    p.font.size = Pt(14)
    p.font.bold = True

    if recommendations:
        for item in recommendations:
            p = tf2.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13)
    else:
        p = tf2.add_paragraph()
        p.text = "None"
        p.level = 1
        p.font.size = Pt(13)

    p = tf2.add_paragraph()
    p.text = ""
    p.level = 0

    p = tf2.add_paragraph()
    p.text = "Action Items:"
    p.level = 0
    p.font.size = Pt(14)
    p.font.bold = True

    if action_items:
        for item in action_items:
            p = tf2.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13)
    else:
        p = tf2.add_paragraph()
        p.text = "None"
        p.level = 1
        p.font.size = Pt(13)

# Save the presentation
prs.save('artifacts/kickoff-deck.pptx')
print(f"\nSuccessfully updated kickoff-deck.pptx with 2 new slides")
print(f"Total slides now: {len(prs.slides)}")
