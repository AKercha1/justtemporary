#!/usr/bin/env python3
"""
Script to add meeting summary slides to the kickoff-deck.pptx presentation.
Adds 2 new slides at the end with meeting summary and action items.
"""

from pptx import Presentation
from pptx.util import Pt

# Load the presentation
prs = Presentation('/app/temp/533734dc-8a59-46fb-b588-cb11695d1b45/artifacts/kickoff-deck.pptx')

# Use the "Content Slide_No Eyebrow" layout (index 5)
content_layout = prs.slide_layouts[5]

# ===== SLIDE 1: Meeting Summary =====
slide1 = prs.slides.add_slide(content_layout)

# Find title and content placeholders
placeholders = list(slide1.placeholders)
if len(placeholders) >= 2:
    # First placeholder is typically the title
    placeholders[0].text = "Meeting Summary: Sheep vs. Goats"

    # Second placeholder is content
    content1 = placeholders[1].text_frame
    content1.clear()

    # Create bullet points for key differentiators
    key_points = [
        ("Dietary Habits", "Sheep are grazers; Goats are browsers (prefer higher vegetation)"),
        ("Physical Features", "Sheep: wool, hanging tails; Goats: hair, upright tails, often beards"),
        ("Behavior", "Sheep are flock animals; Goats are curious explorers"),
        ("Vocalization", "Sheep 'baa'; Goats 'bleat' with raspier, attention-seeking sounds"),
        ("Farm Purpose", "Sheep: wool, meat, milk; Goats: milk, meat, fiber, brush management"),
    ]

    for heading, description in key_points:
        p = content1.add_paragraph()
        p.text = heading
        p.level = 0
        p.font.bold = True
        p.font.size = Pt(14)

        # Add sub-point with description
        p_sub = content1.add_paragraph()
        p_sub.text = description
        p_sub.level = 1
        p_sub.font.size = Pt(12)

# ===== SLIDE 2: Agreed Recommendations & Action Items =====
slide2 = prs.slides.add_slide(content_layout)

# Find title and content placeholders
placeholders2 = list(slide2.placeholders)
if len(placeholders2) >= 2:
    placeholders2[0].text = "Recommendations & Action Items"

    content2 = placeholders2[1].text_frame
    content2.clear()

    recommendations = [
        "Dietary Management: Provide appropriate nutrition, clean water, shelter, and parasite control",
        "Behavioral Awareness: Sheep seek flock safety; goats explore and test boundaries",
        "Movement Differences: Account for goats' climbing ability and sheep's preference for open spaces",
        "Purpose-Driven Selection: Choose based on farm goals",
        "Core Distinction: Sheep graze together; Goats browse and explore",
    ]

    for i, rec in enumerate(recommendations, 1):
        p = content2.add_paragraph()
        p.text = rec
        p.level = 0
        p.font.size = Pt(13)

# Save the updated presentation
prs.save('/app/temp/533734dc-8a59-46fb-b588-cb11695d1b45/artifacts/kickoff-deck.pptx')
print("✓ Successfully updated kickoff-deck.pptx")
print(f"✓ Added 2 new slides (total slides now: {len(prs.slides)})")
