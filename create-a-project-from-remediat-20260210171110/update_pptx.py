#!/usr/bin/env python3
"""
Script to add meeting summary and action items slides to kickoff-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Load the presentation
prs = Presentation('artifacts/kickoff-deck.pptx')

print(f"Total slides in deck: {len(prs.slides)}")

# Find the appropriate layout for content slides
slide_layout = None
for layout in prs.slide_layouts:
    if 'Content Slide_with Eyebrow' in layout.name:
        slide_layout = layout
        break

if not slide_layout:
    # Fallback to the first content layout
    slide_layout = prs.slide_layouts[1]

print(f"Using layout: {slide_layout.name}")

# Extract meeting summary and action items from meeting-summary.md
with open('artifacts/meeting-summary.md', 'r') as f:
    summary_content = f.read()

# Parse the summary
lines = summary_content.split('\n')
meeting_date = None
meeting_time = None
conducted_by = None
purpose = None
key_discussions = []
recommendations = []
action_items = []

current_section = None
for line in lines:
    line = line.strip()
    if 'Date:' in line:
        meeting_date = line.replace('**Date:**', '').strip()
    elif 'Time:' in line:
        meeting_time = line.replace('**Time:**', '').strip()
    elif 'Conducted By:' in line:
        conducted_by = line.replace('**Conducted By:**', '').strip()
    elif 'Purpose:' in line:
        current_section = 'purpose'
    elif 'Key Discussions' in line:
        current_section = 'discussions'
    elif 'Discussed Recommendations' in line:
        current_section = 'recommendations'
    elif 'Agreed Action Items' in line:
        current_section = 'action_items'
    elif line.startswith('*') and current_section == 'discussions':
        key_discussions.append(line.lstrip('* '))
    elif line.startswith('*') and current_section == 'recommendations':
        recommendations.append(line.lstrip('* '))
    elif line.startswith('*') and current_section == 'action_items':
        action_items.append(line.lstrip('* '))
    elif line and current_section == 'purpose' and not line.startswith('**'):
        purpose = line

# --- SLIDE 1: Meeting Summary ---
slide1 = prs.slides.add_slide(slide_layout)

# Find and set title placeholder
title_shape = None
content_shape = None

for shape in slide1.placeholders:
    if shape.placeholder_format.type == 1:  # Title
        title_shape = shape
    elif shape.placeholder_format.type == 2:  # Body
        content_shape = shape

if title_shape:
    title_shape.text = "Meeting Summary"

# Set body content
if content_shape and hasattr(content_shape, 'text_frame'):
    tf = content_shape.text_frame
    tf.clear()

    # Add summary content
    summary_lines = []
    if conducted_by and conducted_by != '[Not specified]':
        summary_lines.append(f"Conducted By: {conducted_by}")
    if meeting_date and meeting_date != '[Date of recording - not specified in VTT]':
        summary_lines.append(f"Date: {meeting_date}")
    if meeting_time and meeting_time != '[Time of recording - not specified in VTT]':
        summary_lines.append(f"Time: {meeting_time}")

    if purpose:
        summary_lines.append("")
        summary_lines.append(f"Purpose: {purpose}")

    if key_discussions:
        summary_lines.append("")
        summary_lines.append("Key Discussion Points:")
        for item in key_discussions:
            summary_lines.append(f"  • {item}")

    for idx, line_text in enumerate(summary_lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.level = 0
        if line_text.startswith('  •'):
            p.level = 1
        p.font.size = Pt(14)

# --- SLIDE 2: Agreed Recommendations & Action Items ---
slide2 = prs.slides.add_slide(slide_layout)

# Find and set title placeholder
title_shape2 = None
content_shape2 = None

for shape in slide2.placeholders:
    if shape.placeholder_format.type == 1:  # Title
        title_shape2 = shape
    elif shape.placeholder_format.type == 2:  # Body
        content_shape2 = shape

if title_shape2:
    title_shape2.text = "Recommendations & Action Items"

# Set body content
if content_shape2 and hasattr(content_shape2, 'text_frame'):
    tf2 = content_shape2.text_frame
    tf2.clear()

    content_lines = []

    if recommendations:
        content_lines.append("Agreed Recommendations:")
        for item in recommendations:
            content_lines.append(f"  • {item}")
    else:
        content_lines.append("Agreed Recommendations:")
        content_lines.append("  • None")

    content_lines.append("")

    if action_items:
        content_lines.append("Action Items:")
        for item in action_items:
            content_lines.append(f"  • {item}")
    else:
        content_lines.append("Action Items:")
        content_lines.append("  • None")

    for idx, line_text in enumerate(content_lines):
        if idx == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line_text
        p.level = 0
        if line_text.startswith('  •'):
            p.level = 1
        p.font.size = Pt(14)

# Save the presentation
prs.save('artifacts/kickoff-deck.pptx')
print(f"Successfully updated kickoff-deck.pptx with 2 new slides")
print(f"Total slides now: {len(prs.slides)}")
